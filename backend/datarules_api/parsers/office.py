from pathlib import Path
from typing import Any

from .common import CanonicalBlock, ParseResult, ParserError


def parse_xlsx(path: Path) -> ParseResult:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise ParserError("openpyxl is required for XLSX parsing") from exc

    workbook = load_workbook(path, data_only=False, read_only=False)
    blocks: list[CanonicalBlock] = []
    for sheet in workbook.worksheets:
        rows: list[list[dict]] = []
        for row in sheet.iter_rows():
            cells = []
            for cell in row:
                if cell.value is None and not cell.data_type:
                    continue
                cells.append(
                    {
                        "cell": cell.coordinate,
                        "value": _safe_value(cell.value),
                        "formula": _safe_value(cell.value) if cell.data_type == "f" else None,
                    }
                )
            if cells:
                rows.append(cells)

        merged = [str(rng) for rng in sheet.merged_cells.ranges]
        blocks.append(
            CanonicalBlock(
                block_type="table",
                sheet_name=sheet.title,
                table_json={"rows": rows, "merged_cells": merged},
                text=f"Sheet {sheet.title}: {len(rows)} non-empty rows",
                confidence=0.98,
            )
        )
    return ParseResult(file_type="xlsx", blocks=blocks, metadata={"sheets": workbook.sheetnames})


def parse_docx(path: Path) -> ParseResult:
    try:
        from docx import Document
    except ImportError as exc:
        raise ParserError("python-docx is required for DOCX parsing") from exc

    document = Document(path)
    blocks: list[CanonicalBlock] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            blocks.append(CanonicalBlock(block_type="paragraph", text=text, confidence=0.98))

    for table in document.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        blocks.append(
            CanonicalBlock(
                block_type="table",
                text="\n".join(" | ".join(row) for row in rows[:30]),
                table_json={"rows": rows},
                confidence=0.95,
            )
        )
    return ParseResult(file_type="docx", blocks=blocks, metadata={"blocks": len(blocks)})


def parse_pptx(path: Path) -> ParseResult:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ParserError("python-pptx is required for PPTX parsing") from exc

    presentation = Presentation(path)
    blocks: list[CanonicalBlock] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                blocks.append(
                    CanonicalBlock(
                        block_type="table",
                        slide_number=idx,
                        text="\n".join(" | ".join(row) for row in rows[:30]),
                        table_json={"rows": rows},
                        confidence=0.92,
                    )
                )
        if text_parts:
            blocks.append(
                CanonicalBlock(
                    block_type="slide",
                    slide_number=idx,
                    text="\n".join(text_parts),
                    confidence=0.94,
                )
            )
    return ParseResult(file_type="pptx", blocks=blocks, metadata={"slides": len(presentation.slides)})


def _safe_value(value: Any) -> Any:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if hasattr(value, "isoformat"):
        return value.isoformat()
    return str(value)
